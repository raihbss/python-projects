import copy
import os
from pptx import Presentation
import pptx


def merge_ppt(paths,final_name):
    outputPres = Presentation(paths[0])
    for path in paths[1:]:
        templatePres = Presentation(path)
        for i in range(len(templatePres.slides)):
            create_new_slide(templatePres, i, outputPres)
    if final_name:
        outputPres.save(final_name)
    return outputPres


def create_new_slide(copyFromPres, slideIndex , pasteIntoPres):
    slide_to_copy = copyFromPres.slides[slideIndex]
    slide_layout = pasteIntoPres.slide_layouts[0]
    new_slide = pasteIntoPres.slides.add_slide(slide_layout)
    imgDict = {}
    for shp in slide_to_copy.shapes:
        if type(shp.element) is pptx.oxml.shapes.picture.CT_Picture:
            with open(shp.name + '.jpg', 'wb') as f:
                f.write(shp.image.blob)
            imgDict[shp.name + '.jpg'] = [shp.left, shp.top, shp.width, shp.height]
        el = shp.element
        newel = copy.deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
    for k, v in imgDict.items():
        new_slide.shapes.add_picture(k, v[0], v[1], v[2], v[3])
        os.remove(k)
    new_slide.shapes.title.text = ' '
    return new_slide





if __name__ == "__main__":
    merge_ppt(paths=["ppt7.pptx", "ppt5.pptx"], final_name="sample.pptx")

